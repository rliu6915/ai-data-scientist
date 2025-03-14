from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData

# Create a presentation object
prs = Presentation()

# Title Slide
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]

title.text = "Sales Analysis Report"
subtitle.text = "Last Quarter Performance and Predictions"

# Slide 1: Top-Performing Category
slide1 = prs.slides.add_slide(prs.slide_layouts[1])
title1 = slide1.shapes.title
content1 = slide1.placeholders[1]

title1.text = "Top-Performing Category"
content1.text = "The top-performing category in terms of sales amount in the last quarter was **Clothing**, with total sales of approximately **$111,051,000**."

# Slide 2: Regression Analysis Results
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]

title2.text = "Predicted Sales"
content2.text = "The regression analysis has been successfully executed, and the predicted sales for the coming quarter (April to June) are:\n - **April**: $9,467,458.48\n - **May**: $9,499,160.77\n - **June**: $9,530,863.07"

# Save the presentation
pptx_file_path = 'output/Sales_Analysis_Report.pptx'
prs.save(pptx_file_path)
